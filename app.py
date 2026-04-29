"""
CHI Insurance Brokers — Αυτόματη Δημιουργία Παρουσίασης
Χρήση: streamlit run app.py
"""

import streamlit as st
import anthropic
import base64
import json
import io
import re
from pathlib import Path
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ─── PAGE CONFIG ────────────────────────────────────────────────────
st.set_page_config(
    page_title="CHI Insurance — Παρουσιάσεις",
    page_icon="🛡️",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ─── COLORS ─────────────────────────────────────────────────────────
C = {
    "navy":      RGBColor(0x1C, 0x3F, 0x5E),
    "navyDark":  RGBColor(0x0F, 0x26, 0x38),
    "teal":      RGBColor(0x00, 0xB4, 0xD8),
    "white":     RGBColor(0xFF, 0xFF, 0xFF),
    "offWhite":  RGBColor(0xF4, 0xF9, 0xFF),
    "textDark":  RGBColor(0x1A, 0x2B, 0x3C),
    "green":     RGBColor(0x27, 0xAE, 0x60),
    "orange":    RGBColor(0xE6, 0x7E, 0x22),
    "red":       RGBColor(0xE7, 0x4C, 0x3C),
    "gold":      RGBColor(0xF5, 0x9E, 0x0B),
    "generali":  RGBColor(0xCC, 0x00, 0x00),
    "morgan":    RGBColor(0x1C, 0x3F, 0x5E),
    "now":       RGBColor(0x7B, 0x2D, 0x8B),
    "blue":      RGBColor(0x3B, 0x82, 0xF6),
}

# ─── HELPERS ────────────────────────────────────────────────────────
def rgb(r, g, b): return RGBColor(r, g, b)

def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_width=None):
    from pptx.util import Pt as _Pt
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width or 1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, size=12, bold=False, italic=False,
             color=None, align=PP_ALIGN.LEFT, valign="middle", bg=None, wrap=True):
    from pptx.util import Pt as _Pt
    from pptx.enum.text import MSO_ANCHOR
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    if valign == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif valign == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txBox

def insurer_color(name):
    n = name.upper()
    if "GENERALI" in n: return C["generali"]
    if "MORGAN" in n:   return C["navy"]
    if "NOW" in n:      return C["now"]
    if "ERGO" in n:     return rgb(0x00, 0x5A, 0xA0)
    return C["teal"]

# ─── CLAUDE PDF EXTRACTION ──────────────────────────────────────────
EXTRACT_PROMPT = """
Διάβασε αυτή την ασφαλιστική προσφορά και εξάγαγε τα παρακάτω στοιχεία σε JSON.
Απάντησε ΜΟΝΟ με valid JSON, χωρίς markdown ή backticks.

{
  "insurer": "Όνομα ασφαλιστικής (π.χ. Generali, Morgan Price, NOW Health, ERGO)",
  "plan_name": "Ακριβές όνομα πλάνου (π.χ. Life On Family Flexi, Evolution Standard Plus)",
  "annual_premium": "Ετήσιο ασφάλιστρο σε αριθμό (μόνο νούμερο, π.χ. 2626)",
  "currency": "EUR ή USD ή GBP",
  "deductible": "Απαλλαγή (π.χ. 500, ή 1000 ανά άτομο, ή 1500 κοινή)",
  "max_coverage": "Μέγιστο κεφάλαιο (π.χ. 750000)",
  "geography": "Γεωγραφική κάλυψη (π.χ. Ευρώπη, Παγκόσμια εκτός ΗΠΑ)",
  "hospital_class": "Θέση νοσηλείας (π.χ. A, B, Standard Room)",
  "inpatient": "Full Refund ή ποσοστό ή Not Covered",
  "outpatient_limit": "Όριο εξωνοσοκομειακών (π.χ. 2500 ή Not Covered)",
  "outpatient_pct": "Ποσοστό κάλυψης εξωνοσοκομειακών (π.χ. 80 ή 100 ή null)",
  "mri_ct_pet": "Full Refund ή Not Covered ή περιγραφή",
  "cancer": "Full Refund ή περιγραφή ή Not Covered",
  "physiotherapy": "Ποσό ή Full Refund ή Not Covered",
  "psychiatric_outpatient": "Περιγραφή ή Not Covered",
  "home_nursing": "Περιγραφή ή Not Covered",
  "waiting_period": "Αναμονή για παθήσεις (π.χ. Άμεση ή 6 μήνες ή 24 μήνες)",
  "preexisting": "Κάλυψη προϋπαρχουσών (π.χ. Άμεση MHD ή μετά 12μήνες ή Όχι)",
  "key_notes": ["σημαντική παρατήρηση 1", "σημαντική παρατήρηση 2"],
  "insured_members": [
    {"age": 54, "role": "Κύρια Ασφαλισμένη"},
    {"age": 17, "role": "Εξαρτώμενο Μέλος"}
  ]
}

Αν κάποιο πεδίο δεν βρεθεί, βάλε null.
"""

def extract_insurance_data(pdf_bytes: bytes, api_key: str) -> dict:
    client = anthropic.Anthropic(api_key=api_key)
    pdf_b64 = base64.standard_b64encode(pdf_bytes).decode("utf-8")
    
    response = client.messages.create(
        model="claude-opus-4-5",
        max_tokens=2000,
        messages=[{
            "role": "user",
            "content": [
                {
                    "type": "document",
                    "source": {
                        "type": "base64",
                        "media_type": "application/pdf",
                        "data": pdf_b64,
                    },
                },
                {"type": "text", "text": EXTRACT_PROMPT}
            ],
        }]
    )
    
    raw = response.content[0].text.strip()
    # Clean any accidental markdown
    raw = re.sub(r"```json|```", "", raw).strip()
    return json.loads(raw)


# ─── PPTX GENERATION ────────────────────────────────────────────────
def generate_pptx(client_name: str, client_members: list, proposals: list,
                  recommended_idx: int, broker_name: str, broker_tel: str, broker_email: str) -> bytes:
    """
    proposals: list of dicts with extracted insurance data
    recommended_idx: index of recommended proposal (0-based)
    """
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank
    footer_text = f"CHI Insurance Brokers | {broker_name}  ·  {broker_tel}  ·  {broker_email}"

    def add_footer(slide, light=False):
        fc = C["offWhite"] if not light else C["navyDark"]
        add_rect(slide, 0, 7.18, 13.33, 0.32, C["navy"] if light else C["navyDark"])
        add_text(slide, footer_text, 0.3, 7.18, 12.7, 0.32,
                 size=8, color=C["teal"], align=PP_ALIGN.CENTER)

    def add_top_bar(slide, color):
        add_rect(slide, 0, 0, 13.33, 0.1, color)

    # ── SLIDE 1: COVER ──────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C["navyDark"]

    add_top_bar(s, C["teal"])
    add_rect(s, 0, 0.1, 5.0, 7.08, C["navy"])
    add_rect(s, 5.0, 0.1, 0.05, 7.08, C["teal"])

    add_text(s, "ΑΣΦΑΛΙΣΗ ΥΓΕΙΑΣ", 0.3, 2.8, 4.4, 0.5,
             size=12, bold=True, color=C["teal"], align=PP_ALIGN.CENTER)
    add_text(s, "ΣΥΓΚΡΙΤΙΚΗ ΑΝΑΛΥΣΗ", 0.3, 3.25, 4.4, 0.4,
             size=10, color=rgb(0x9D, 0xC4, 0xD8), align=PP_ALIGN.CENTER)

    add_text(s, "Πρόταση", 5.5, 1.0, 7.5, 0.8,
             size=22, italic=True, color=C["teal"])
    add_text(s, "Ασφάλισης Υγείας", 5.5, 1.7, 7.5, 1.2,
             size=44, bold=True, color=C["white"])
    add_rect(s, 5.5, 3.0, 7.6, 0.05, C["teal"])

    # Client names
    names_text = client_name
    add_text(s, names_text, 5.5, 3.15, 7.5, 0.6,
             size=20, bold=True, color=C["white"])
    member_str = "  ·  ".join([f"{m.get('role','')} ({m.get('age','')} ετών)" for m in client_members])
    add_text(s, member_str, 5.5, 3.75, 7.5, 0.45,
             size=12, color=rgb(0x9D, 0xC4, 0xD8))

    add_text(s, datetime.now().strftime("%B %Y"), 5.5, 4.4, 2.0, 0.45,
             size=11, bold=True, color=C["navy"], align=PP_ALIGN.CENTER,
             bg=C["teal"])
    add_rect(s, 5.5, 4.4, 2.0, 0.45, C["teal"])
    add_text(s, datetime.now().strftime("%B %Y"), 5.5, 4.4, 2.0, 0.45,
             size=11, bold=True, color=C["navy"], align=PP_ALIGN.CENTER)

    # Insurer pills
    for i, prop in enumerate(proposals[:4]):
        xp = 5.5 + i * 1.95
        col = insurer_color(prop.get("insurer", ""))
        add_rect(s, xp, 5.2, 1.8, 0.42, col)
        add_text(s, prop.get("insurer", "").upper(), xp, 5.2, 1.8, 0.42,
                 size=8, bold=True, color=C["white"], align=PP_ALIGN.CENTER)

    add_footer(s, light=True)

    # ── SLIDE 2: OVERVIEW ───────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C["navy"]
    add_top_bar(s, C["gold"])

    add_text(s, f"Τρεις Προτάσεις για {client_name}", 0.4, 0.2, 12.5, 0.65,
             size=28, bold=True, color=C["white"])
    add_text(s, "Από την οικονομική επιλογή έως την ολοκληρωμένη διεθνή κάλυψη",
             0.4, 0.85, 12.5, 0.4, size=13, italic=True, color=rgb(0x9D,0xC4,0xD8))

    n = len(proposals)
    card_w = (13.33 - 0.4 - 0.4 - (n-1)*0.25) / n
    for i, prop in enumerate(proposals):
        xp = 0.4 + i * (card_w + 0.25)
        is_rec = (i == recommended_idx)
        col = insurer_color(prop.get("insurer",""))
        bg_col = rgb(0x1A,0x2F,0x45) if is_rec else rgb(0x17,0x35,0x4E)
        border = C["gold"] if is_rec else col

        add_rect(s, xp, 1.4, card_w, 5.5, bg_col, line_color=border, line_width=2 if is_rec else 0.8)
        add_rect(s, xp, 1.4, card_w, 0.5, col)

        label = f"{'★ ΠΡΟΤΕΙΝΟΜΕΝΗ' if is_rec else f'ΕΠΙΛΟΓΗ {chr(65+i)}'}"
        add_text(s, label, xp, 1.4, card_w, 0.5,
                 size=9, bold=True, color=C["navy"] if is_rec else C["white"], align=PP_ALIGN.CENTER)

        add_text(s, prop.get("insurer",""), xp+0.1, 2.0, card_w-0.2, 0.42,
                 size=11, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_text(s, prop.get("plan_name",""), xp+0.1, 2.42, card_w-0.2, 0.65,
                 size=12, bold=True, color=C["white"], align=PP_ALIGN.CENTER)

        price_bg = C["gold"] if is_rec else rgb(0x2A,0x4A,0x63)
        price_fg = C["navy"] if is_rec else C["teal"]
        add_rect(s, xp+card_w*0.15, 3.2, card_w*0.7, 0.65, price_bg)
        cur = prop.get("currency","€") if prop.get("currency") else "€"
        sym = "€" if cur=="EUR" else ("$" if cur=="USD" else "£")
        premium = prop.get("annual_premium","")
        add_text(s, f"{sym}{premium}" if premium else "—", xp+card_w*0.15, 3.2, card_w*0.7, 0.65,
                 size=20, bold=True, color=price_fg, align=PP_ALIGN.CENTER)

        add_text(s, prop.get("geography",""), xp+0.1, 4.0, card_w-0.2, 0.35,
                 size=9, italic=True, color=rgb(0x9D,0xC4,0xD8), align=PP_ALIGN.CENTER)

        # Mini bullets
        bullets = [
            f"Απαλλαγή: {prop.get('deductible','—')}",
            f"Νοσηλεία: {prop.get('inpatient','—')}",
            f"MRI/PET: {prop.get('mri_ct_pet','—')}",
            f"Αναμονή: {prop.get('waiting_period','—')}",
        ]
        for bi, b in enumerate(bullets):
            add_text(s, f"• {b}", xp+0.15, 4.45+bi*0.47, card_w-0.3, 0.42,
                     size=9, color=rgb(0xC8,0xDF,0xF0))

    add_footer(s, light=True)

    # ── SLIDES 3+: ONE PER PROPOSAL ─────────────────────────────────
    for i, prop in enumerate(proposals):
        s = prs.slides.add_slide(blank_layout)
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = C["offWhite"]

        col = insurer_color(prop.get("insurer",""))
        is_rec = (i == recommended_idx)
        if is_rec:
            s.background.fill.fore_color.rgb = C["navyDark"]

        add_top_bar(s, C["gold"] if is_rec else col)
        add_rect(s, 0, 0.1, 0.35, 7.08, C["gold"] if is_rec else col)

        label = f"{'★ ΕΠΙΛΟΓΗ ' + chr(65+i) + ' — ΠΡΟΤΕΙΝΟΜΕΝΗ' if is_rec else 'ΕΠΙΛΟΓΗ ' + chr(65+i)}"
        add_text(s, label, 0.5, 0.15, 7.0, 0.42,
                 size=10, bold=True, color=C["gold"] if is_rec else col)
        add_text(s, f"{prop.get('insurer','')} — {prop.get('plan_name','')}", 0.5, 0.55, 10.0, 0.75,
                 size=26, bold=True, color=C["white"] if is_rec else C["navy"])

        # Price badge
        price_col = C["gold"] if is_rec else col
        add_rect(s, 10.2, 0.3, 2.8, 1.0, price_col)
        cur = prop.get("currency","€") or "€"
        sym = "€" if cur=="EUR" else ("$" if cur=="USD" else "£")
        add_text(s, f"{sym}{prop.get('annual_premium','—')}", 10.2, 0.3, 2.8, 0.65,
                 size=28, bold=True, color=C["navy"] if is_rec else C["white"], align=PP_ALIGN.CENTER)
        add_text(s, "ετήσιο ασφάλιστρο", 10.2, 0.93, 2.8, 0.37,
                 size=9, color=rgb(0xC8,0xDF,0xF0) if not is_rec else rgb(0x5A,0x4A,0x00), align=PP_ALIGN.CENTER)

        # Key params
        params = [
            ("Μέγιστο Κεφάλαιο", f"{sym}{prop.get('max_coverage','—')}"),
            ("Απαλλαγή", prop.get("deductible","—") or "—"),
            ("Γεωγραφία", prop.get("geography","—") or "—"),
            ("Θέση Νοσηλείας", prop.get("hospital_class","—") or "—"),
        ]
        for pi, (k, v) in enumerate(params):
            yp = 1.55 + pi * 0.52
            add_rect(s, 0.5, yp, 3.2, 0.42, C["navy"] if not is_rec else rgb(0x1A,0x35,0x50))
            add_text(s, k, 0.5, yp, 3.2, 0.42,
                     size=10, bold=True, color=C["white"] if not is_rec else C["gold"], align=PP_ALIGN.CENTER)
            add_text(s, v, 3.8, yp+0.05, 9.0, 0.42,
                     size=10, color=C["textDark"] if not is_rec else rgb(0xC8,0xDF,0xF0))

        # Coverage grid
        coverage_col = C["green"]
        add_rect(s, 0.5, 3.75, 6.0, 0.4, coverage_col)
        add_text(s, "✓  ΚΑΛΥΨΕΙΣ", 0.5, 3.75, 6.0, 0.4,
                 size=10, bold=True, color=C["white"], align=PP_ALIGN.CENTER)

        covers = [
            ("Νοσηλεία", prop.get("inpatient","—") or "—"),
            ("Εξωνοσοκ.", f"εως {sym}{prop.get('outpatient_limit','—')} ({prop.get('outpatient_pct','—')}%)" if prop.get("outpatient_limit") and prop.get("outpatient_limit") != "Not Covered" else prop.get("outpatient_limit","—") or "—"),
            ("MRI/CT/PET", prop.get("mri_ct_pet","—") or "—"),
            ("Καρκίνος", prop.get("cancer","—") or "—"),
            ("Φυσιοθεραπεία", prop.get("physiotherapy","—") or "—"),
        ]
        for ci, (k, v) in enumerate(covers):
            yp = 4.25 + ci * 0.47
            tick = "✓" if "Not Covered" not in str(v) and v != "—" else "✗"
            tc = C["green"] if tick == "✓" else C["red"]
            add_text(s, tick, 0.55, yp, 0.4, 0.42, size=12, bold=True, color=tc, align=PP_ALIGN.CENTER)
            add_text(s, k, 1.0, yp, 1.5, 0.42, size=10, bold=True, color=C["white"] if is_rec else C["navy"])
            add_text(s, v, 2.6, yp, 3.7, 0.42, size=9, color=rgb(0xC8,0xDF,0xF0) if is_rec else C["textDark"])

        # Notes / Limitations
        add_rect(s, 6.8, 3.75, 6.2, 0.4, C["orange"])
        add_text(s, "⚠  ΣΗΜΑΝΤΙΚΕΣ ΠΑΡΑΤΗΡΗΣΕΙΣ", 6.8, 3.75, 6.2, 0.4,
                 size=10, bold=True, color=C["white"], align=PP_ALIGN.CENTER)

        notes = prop.get("key_notes", [])
        notes_extra = [
            f"Αναμονή: {prop.get('waiting_period','Άμεση')}",
            f"Προϋπ. παθήσεις: {prop.get('preexisting','—')}",
        ]
        all_notes = (notes or []) + notes_extra
        for ni, note in enumerate(all_notes[:6]):
            yp = 4.25 + ni * 0.5
            add_text(s, f"• {note}", 6.85, yp, 6.1, 0.45,
                     size=9, color=rgb(0xC8,0xDF,0xF0) if is_rec else C["textDark"])

        add_footer(s, light=is_rec)

    # ── COMPARISON TABLE SLIDE ──────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C["offWhite"]
    add_top_bar(s, C["teal"])

    add_text(s, "Σύγκριση Καλύψεων — Πίνακας", 0.4, 0.2, 12.5, 0.6,
             size=26, bold=True, color=C["navy"])

    # Table header
    headers = ["ΚΑΛΥΨΗ"] + [f"{p.get('insurer','')}\n{p.get('plan_name','')}" for p in proposals]
    col_widths = [3.5] + [(13.33 - 0.4 - 0.4 - 3.5) / len(proposals)] * len(proposals)
    xstarts = [0.4]
    for cw in col_widths[:-1]:
        xstarts.append(xstarts[-1] + cw)

    header_colors = [C["navy"]] + [insurer_color(p.get("insurer","")) for p in proposals]
    header_text_colors = [C["white"]] * len(headers)
    # Gold for recommended
    if recommended_idx < len(proposals):
        header_colors[recommended_idx + 1] = C["gold"]
        header_text_colors[recommended_idx + 1] = C["navy"]

    for ci, (h, xp, cw, hc, htc) in enumerate(zip(headers, xstarts, col_widths, header_colors, header_text_colors)):
        add_rect(s, xp, 0.95, cw, 0.6, hc)
        add_text(s, h, xp+0.05, 0.95, cw-0.1, 0.6,
                 size=8, bold=True, color=htc, align=PP_ALIGN.CENTER)

    # Table rows
    row_data = [
        ("Max Κεφάλαιο/Έτος", lambda p: f"€{p.get('max_coverage','—')}"),
        ("Απαλλαγή", lambda p: p.get("deductible","—") or "—"),
        ("Γεωγραφία", lambda p: p.get("geography","—") or "—"),
        ("Θέση Νοσηλείας", lambda p: p.get("hospital_class","—") or "—"),
        ("Νοσηλεία In-Patient", lambda p: p.get("inpatient","—") or "—"),
        ("MRI/CT/PET (εξωτ.)", lambda p: p.get("mri_ct_pet","—") or "—"),
        ("Εξωνοσ. Ιατροί/Εξετ.", lambda p: f"εως €{p.get('outpatient_limit','—')} ({p.get('outpatient_pct','—')}%)" if p.get("outpatient_limit") and p.get("outpatient_limit") != "Not Covered" else (p.get("outpatient_limit","—") or "—")),
        ("Φυσικοθεραπεία", lambda p: p.get("physiotherapy","—") or "—"),
        ("Καρκίνος", lambda p: p.get("cancer","—") or "—"),
        ("Ψυχιατρική Εξωτ.", lambda p: p.get("psychiatric_outpatient","—") or "—"),
        ("Αναμονή Κάλυψης", lambda p: p.get("waiting_period","—") or "—"),
        ("Ετήσιο Ασφάλιστρο", lambda p: f"€{p.get('annual_premium','—')}"),
    ]

    for ri, (label, fn) in enumerate(row_data):
        yp = 1.65 + ri * 0.43
        row_bg = rgb(0xF4,0xF9,0xFF) if ri % 2 == 0 else C["white"]
        for ci, (xp, cw) in enumerate(zip(xstarts, col_widths)):
            bg = row_bg
            is_rec_col = ci == recommended_idx + 1
            if is_rec_col:
                bg = rgb(0xFF,0xF8,0xE1) if ri % 2 == 0 else rgb(0xFF,0xF3,0xC4)
            add_rect(s, xp, yp, cw, 0.41, bg)
            if ci == 0:
                add_text(s, label, xp+0.08, yp, cw-0.1, 0.41, size=9, bold=False, color=C["navy"])
            else:
                val = fn(proposals[ci-1])
                is_good = val not in ["Not Covered", "—", "null", None] and "Not Covered" not in str(val)
                col_t = C["green"] if is_good else C["red"]
                ft = val if val not in ["Not Covered", "null", None] else "ΔΕΝ καλύπτει"
                add_text(s, str(ft), xp+0.05, yp, cw-0.1, 0.41,
                         size=8.5, bold=is_rec_col, color=col_t if ci > 0 and "€" not in str(ft) else (C["navy"] if is_rec_col else C["navy"]),
                         align=PP_ALIGN.CENTER)

    add_footer(s)

    # ── CLOSING SLIDE ───────────────────────────────────────────────
    s = prs.slides.add_slide(blank_layout)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C["navyDark"]
    add_top_bar(s, C["gold"])

    add_text(s, "Η Πρότασή μας", 0.5, 0.25, 12.0, 0.7,
             size=30, bold=True, color=C["gold"])

    if recommended_idx < len(proposals):
        rec = proposals[recommended_idx]
        rec_text = f"{rec.get('insurer','')} — {rec.get('plan_name','')}  |  €{rec.get('annual_premium','')} / έτος"
        add_text(s, rec_text, 0.5, 1.0, 12.5, 0.6, size=16, bold=True, color=C["white"])

    steps = [
        ("ΒΗΜΑ 1", "Εντός 48ωρών", "Έγκριση πρότασης & αποστολή ιατρικού ιστορικού"),
        ("ΒΗΜΑ 2", "Underwriting", "Υπογραφή — σαφής γνώση τι καλύπτεται"),
        ("ΒΗΜΑ 3", "Ενεργοποίηση", "Άμεση κάλυψη χωρίς αναμονές"),
    ]
    for si, (tag, title, body) in enumerate(steps):
        xp = 0.5 + si * 4.25
        add_rect(s, xp, 1.85, 4.0, 2.6, rgb(0x1A,0x35,0x50))
        add_rect(s, xp, 1.85, 4.0, 0.58, C["gold"])
        add_text(s, tag, xp, 1.85, 4.0, 0.58, size=12, bold=True, color=C["navy"], align=PP_ALIGN.CENTER)
        add_text(s, title, xp+0.1, 2.55, 3.8, 0.55, size=14, bold=True, color=C["teal"], align=PP_ALIGN.CENTER)
        add_text(s, body, xp+0.15, 3.15, 3.7, 1.15, size=10, color=rgb(0xC8,0xDF,0xF0), align=PP_ALIGN.CENTER)

    add_rect(s, 0.5, 4.65, 12.33, 1.05, rgb(0x1A,0x35,0x50))
    add_rect(s, 0.5, 4.65, 0.06, 1.05, C["gold"])
    quote = "Η ασφάλεια υγείας δεν είναι κόστος — είναι επένδυση στην ηρεμία σας και στην οικογένειά σας."
    add_text(s, quote, 0.7, 4.65, 12.0, 1.05, size=13, italic=True, color=rgb(0xC8,0xDF,0xF0))

    add_footer(s, light=True)

    # Save to bytes
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# ─── STREAMLIT UI ───────────────────────────────────────────────────
def main():
    st.markdown("""
    <style>
    .main { background: #F4F9FF; }
    .stButton > button { background: #1C3F5E; color: white; border-radius: 8px;
                          font-weight: bold; padding: 0.6em 2em; border: none; }
    .stButton > button:hover { background: #00B4D8; }
    div[data-testid="stFileUploader"] { border: 2px dashed #00B4D8; border-radius: 8px; padding: 1em; }
    .metric-card { background: white; border-radius: 10px; padding: 1em;
                   border-left: 4px solid #00B4D8; margin-bottom: 0.5em; }
    </style>
    """, unsafe_allow_html=True)

    # Header
    col1, col2 = st.columns([1, 4])
    with col1:
        st.markdown("# 🛡️")
    with col2:
        st.markdown("## CHI Insurance Brokers")
        st.markdown("*Αυτόματη Δημιουργία Παρουσιάσεων Ασφάλισης*")
    st.divider()

    # ── SIDEBAR ──────────────────────────────────────────────────────
    with st.sidebar:
        st.markdown("### ⚙️ Ρυθμίσεις")

        api_key = st.text_input("🔑 Claude API Key", type="password",
                                help="Το API key σου από το console.anthropic.com")

        st.markdown("---")
        st.markdown("### 👤 Στοιχεία Μεσίτη")
        broker_name  = st.text_input("Όνομα", value="Ιατρόπουλος Χρήστος")
        broker_tel   = st.text_input("Τηλέφωνο", value="+30 697 590 0189")
        broker_email = st.text_input("Email", value="info@chiinsurancebrokers.com")

        st.markdown("---")
        st.markdown("### 👥 Στοιχεία Πελάτη")
        client_name = st.text_input("Επώνυμο / Όνομα Πελάτη", placeholder="π.χ. Τοτικίδη Κατία")

        st.markdown("**Μέλη:**")
        n_members = st.number_input("Αριθμός μελών", 1, 6, 2)
        members = []
        for i in range(n_members):
            c1, c2 = st.columns(2)
            with c1:
                age = st.number_input(f"Ηλικία #{i+1}", 0, 99, 30 if i==0 else 17, key=f"age_{i}")
            with c2:
                role = st.selectbox("Ρόλος", ["Κύρια Ασφαλισμένη","Κύριος Ασφαλισμένος","Εξαρτώμενο Μέλος","Σύζυγος"], key=f"role_{i}")
            members.append({"age": age, "role": role})

    # ── MAIN AREA ────────────────────────────────────────────────────
    st.markdown("### 📄 Φόρτωσε τις Ασφαλιστικές Προσφορές (PDF)")
    st.info("Φόρτωσε 2–4 PDF προσφορές. Το Claude θα εξάγει αυτόματα όλα τα στοιχεία.", icon="ℹ️")

    uploaded_files = st.file_uploader(
        "Επίλεξε PDF αρχεία", type="pdf",
        accept_multiple_files=True,
        help="Ανέβασε τις προσφορές Generali, Morgan Price, NOW Health κ.λπ."
    )

    if not uploaded_files:
        st.markdown("---")
        st.markdown("#### Πώς λειτουργεί:")
        c1, c2, c3 = st.columns(3)
        with c1:
            st.markdown("**1️⃣ Ανέβασε PDFs**\nΌλες οι προσφορές που θέλεις να συγκρίνεις")
        with c2:
            st.markdown("**2️⃣ Claude τα αναλύει**\nΕξάγει αυτόματα κεφάλαια, απαλλαγές, καλύψεις")
        with c3:
            st.markdown("**3️⃣ Download PPTX**\nΈτοιμη παρουσίαση με το brand σου")
        return

    # ── EXTRACTION ──────────────────────────────────────────────────
    if "proposals" not in st.session_state:
        st.session_state.proposals = {}

    if st.button("🤖 Ανάλυση με Claude API", type="primary", disabled=not api_key):
        if not api_key:
            st.error("Χρειάζεσαι Claude API key!")
            return

        progress = st.progress(0, text="Αρχικοποίηση...")
        st.session_state.proposals = {}

        for idx, uf in enumerate(uploaded_files):
            progress.progress((idx) / len(uploaded_files),
                              text=f"Ανάλυση: {uf.name}...")
            try:
                data = extract_insurance_data(uf.read(), api_key)
                st.session_state.proposals[uf.name] = data
                st.success(f"✅ {uf.name} → {data.get('insurer','')} {data.get('plan_name','')}")
            except Exception as e:
                st.error(f"❌ Σφάλμα στο {uf.name}: {e}")

        progress.progress(1.0, text="Ολοκληρώθηκε!")

    # ── SHOW EXTRACTED DATA ─────────────────────────────────────────
    if st.session_state.get("proposals"):
        proposals_list = list(st.session_state.proposals.values())
        file_names = list(st.session_state.proposals.keys())

        st.markdown("---")
        st.markdown("### 📊 Εξαχθέντα Στοιχεία")

        # Allow manual edits
        edited_proposals = []
        tabs = st.tabs([f"📋 {p.get('insurer','?')} — {p.get('plan_name','?')[:20]}" for p in proposals_list])
        for tab, prop, fname in zip(tabs, proposals_list, file_names):
            with tab:
                col1, col2 = st.columns(2)
                with col1:
                    prop["insurer"]       = st.text_input("Ασφαλιστική", prop.get("insurer",""), key=f"ins_{fname}")
                    prop["plan_name"]     = st.text_input("Πλάνο", prop.get("plan_name",""), key=f"plan_{fname}")
                    prop["annual_premium"]= st.text_input("Ετήσιο Ασφάλιστρο", str(prop.get("annual_premium","")), key=f"prem_{fname}")
                    prop["deductible"]    = st.text_input("Απαλλαγή", prop.get("deductible",""), key=f"ded_{fname}")
                    prop["geography"]     = st.text_input("Γεωγραφία", prop.get("geography",""), key=f"geo_{fname}")
                    prop["hospital_class"]= st.text_input("Θέση Νοσηλείας", prop.get("hospital_class",""), key=f"hosp_{fname}")
                with col2:
                    prop["inpatient"]     = st.text_input("Νοσηλεία", prop.get("inpatient",""), key=f"inp_{fname}")
                    prop["mri_ct_pet"]    = st.text_input("MRI/CT/PET", prop.get("mri_ct_pet",""), key=f"mri_{fname}")
                    prop["outpatient_limit"] = st.text_input("Εξωνοσοκ. Όριο", str(prop.get("outpatient_limit","")), key=f"outp_{fname}")
                    prop["outpatient_pct"]   = st.text_input("Εξωνοσοκ. %", str(prop.get("outpatient_pct","")), key=f"outpct_{fname}")
                    prop["cancer"]        = st.text_input("Καρκίνος", prop.get("cancer",""), key=f"can_{fname}")
                    prop["waiting_period"]= st.text_input("Αναμονή", prop.get("waiting_period",""), key=f"wait_{fname}")
                edited_proposals.append(prop)

        # Recommended choice
        st.markdown("---")
        st.markdown("### 🎯 Επιλογή Πρότασης")
        insurer_labels = [f"{p.get('insurer','')} — {p.get('plan_name','')} (€{p.get('annual_premium','')})"
                          for p in edited_proposals]
        rec_idx = st.selectbox("Ποια πρόταση να εμφανίζεται ως **ΠΡΟΤΕΙΝΟΜΕΝΗ**?",
                                range(len(insurer_labels)),
                                format_func=lambda i: insurer_labels[i])

        # Generate button
        st.markdown("---")
        if st.button("🎨 Δημιουργία Παρουσίασης PPTX", type="primary"):
            if not client_name:
                st.warning("Συμπλήρωσε το όνομα του πελάτη στο sidebar!")
                return

            with st.spinner("Δημιουργία παρουσίασης..."):
                try:
                    pptx_bytes = generate_pptx(
                        client_name=client_name,
                        client_members=members,
                        proposals=edited_proposals,
                        recommended_idx=rec_idx,
                        broker_name=broker_name,
                        broker_tel=broker_tel,
                        broker_email=broker_email,
                    )

                    fname_out = f"{client_name.replace(' ','_')}_Insurance_{datetime.now().strftime('%Y%m')}.pptx"
                    st.download_button(
                        label="⬇️ Download Παρουσίαση",
                        data=pptx_bytes,
                        file_name=fname_out,
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    )
                    st.success(f"✅ Η παρουσίαση '{fname_out}' είναι έτοιμη!")

                except Exception as e:
                    st.error(f"Σφάλμα: {e}")
                    import traceback
                    st.code(traceback.format_exc())

    elif uploaded_files and not st.session_state.get("proposals"):
        st.info("👆 Πάτα 'Ανάλυση με Claude API' για να εξαχθούν τα στοιχεία από τα PDFs.")


if __name__ == "__main__":
    main()
